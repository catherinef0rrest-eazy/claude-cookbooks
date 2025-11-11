#!/usr/bin/env python3
"""
GTM Fabric PowerPoint Generator
Creates modern, sleek presentations using the ICP Cards design system with GTM Fabric branding.
"""

import os
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


@dataclass
class GTMFabricColors:
    """GTM Fabric ICP Cards Design System Colors"""

    # Background colors (dark theme)
    background: RGBColor = RGBColor(6, 4, 15)  # #06040f - Deep purple-black
    card: RGBColor = RGBColor(18, 10, 31)  # #120a1f - Elevated card
    sidebar: RGBColor = RGBColor(14, 8, 23)  # #0e0817 - Sidebar

    # Text colors
    text_primary: RGBColor = RGBColor(242, 242, 242)  # #f2f2f2 - Main text
    text_secondary: RGBColor = RGBColor(166, 166, 166)  # #a6a6a6 - Secondary text

    # Border colors
    border_default: RGBColor = RGBColor(58, 45, 79)  # #3a2d4f
    border_card: RGBColor = RGBColor(38, 27, 53)  # #261b35

    # Accent color
    accent: RGBColor = RGBColor(192, 132, 252)  # #c084fc - Light Purple

    # Additional utility colors
    white: RGBColor = RGBColor(255, 255, 255)
    black: RGBColor = RGBColor(0, 0, 0)


class GTMFabricPresentationGenerator:
    """Generate modern, sleek PowerPoint presentations with GTM Fabric branding"""

    def __init__(self, logo_path: Optional[str] = None):
        """
        Initialize presentation generator

        Args:
            logo_path: Path to GTM Fabric logo (uses white logo by default for dark theme)
        """
        self.colors = GTMFabricColors()
        self.font_family = "Arial"
        self.font_fallback = ["Calibri", "sans-serif"]

        # Set logo path
        if logo_path is None:
            # Default to white logo for dark theme
            script_dir = os.path.dirname(os.path.abspath(__file__))
            self.logo_path = os.path.join(script_dir, "gtmfabric_logo_white.png")
        else:
            self.logo_path = logo_path

        # Verify logo exists
        if not os.path.exists(self.logo_path):
            print(f"Warning: Logo not found at {self.logo_path}")
            self.logo_path = None

    def create_presentation(self, slide_width: float = 10, slide_height: float = 5.625) -> Presentation:
        """
        Create a new presentation with standard 16:9 dimensions

        Args:
            slide_width: Width in inches (default 10 for standard 16:9)
            slide_height: Height in inches (default 5.625 for standard 16:9)

        Returns:
            Presentation object
        """
        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)
        return prs

    def _apply_dark_background(self, slide):
        """Apply dark theme background to slide"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors.background

    def _add_logo(self, slide, position: str = "top-right", size: float = 0.8,
                  align_center_y: Optional[float] = None):
        """
        Add GTM Fabric logo to slide

        Args:
            slide: PowerPoint slide object
            position: Logo position ("top-right", "top-left", "bottom-right", "bottom-left")
            size: Logo size in inches
            align_center_y: If provided (in inches), centers the logo vertically at this position
        """
        if not self.logo_path:
            return

        if align_center_y is not None:
            # Center logo vertically at align_center_y
            left = Inches(9.3 - size)
            # Logo actual height is size * (actual_aspect_ratio), but we'll approximate
            # For the white logo (6324x1374), aspect = 1374/6324 ≈ 0.217
            logo_aspect = 0.217
            logo_height = size * logo_aspect
            top = Inches(align_center_y) - Inches(logo_height / 2)
        else:
            # Calculate logo position (for 16:9 - 10" x 5.625")
            positions = {
                "top-right": (Inches(9.3 - size), Inches(0.2)),
                "top-left": (Inches(0.3), Inches(0.2)),
                "bottom-right": (Inches(9.3 - size), Inches(5.625 - size - 0.2)),
                "bottom-left": (Inches(0.3), Inches(5.625 - size - 0.2)),
            }
            left, top = positions.get(position, positions["top-right"])

        try:
            slide.shapes.add_picture(
                self.logo_path,
                left,
                top,
                width=Inches(size)
            )
        except Exception as e:
            print(f"Warning: Could not add logo: {e}")

    def _add_footer(self, slide, text: str, page_number: Optional[int] = None):
        """Add footer with optional page number"""
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2),
            Inches(9), Inches(0.3)
        )
        text_frame = footer_box.text_frame
        text_frame.text = text
        if page_number:
            text_frame.text += f" | Page {page_number}"

        para = text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.color.rgb = self.colors.text_secondary
        para.font.name = self.font_family
        para.alignment = PP_ALIGN.CENTER

    def create_title_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: str,
        date: Optional[str] = None,
        include_logo: bool = True
    ):
        """
        Create modern startup UI title slide with geometric grid background

        Args:
            prs: Presentation object
            title: Main title text
            subtitle: Subtitle text
            date: Optional date string
            include_logo: Whether to include GTM Fabric logo
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Apply dark gradient background
        from pptx.enum.dml import MSO_FILL
        from pptx.util import Pt as PtUtil
        from pptx.dml.color import RGBColor

        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135.0  # Diagonal gradient for modern look

        # Multi-stop gradient for depth
        fill.gradient_stops[0].color.rgb = RGBColor(0x0a, 0x08, 0x15)  # Darker purple-black
        fill.gradient_stops[1].color.rgb = self.colors.background  # Base background

        # === GEOMETRIC GRID OVERLAY ===
        # Create subtle grid pattern with lines
        grid_color = RGBColor(0x1a, 0x15, 0x25)  # Subtle purple-gray
        grid_spacing = Inches(1.0)  # 1 inch grid

        # Vertical grid lines
        for i in range(11):  # 10 inches wide
            x_pos = Inches(i * 1.0)
            line = slide.shapes.add_connector(
                1, x_pos, Inches(0), x_pos, Inches(5.625)
            )
            line.line.color.rgb = grid_color
            line.line.width = Pt(0.5)
            line.line.transparency = 0.7  # Make it subtle

        # Horizontal grid lines
        for i in range(6):  # ~5.625 inches tall
            y_pos = Inches(i * 1.0)
            line = slide.shapes.add_connector(
                1, Inches(0), y_pos, Inches(10), y_pos
            )
            line.line.color.rgb = grid_color
            line.line.width = Pt(0.5)
            line.line.transparency = 0.7  # Make it subtle

        # === ACCENT ELEMENTS ===
        # Modern accent line on left
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.3), Inches(2.0),
            Inches(0.08), Inches(2.0)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.colors.accent
        accent_bar.line.fill.background()

        # Accent glow effect (simulated with transparent rectangles)
        glow_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25), Inches(2.0),
            Inches(0.18), Inches(2.0)
        )
        glow_bar.fill.solid()
        glow_bar.fill.fore_color.rgb = self.colors.accent
        glow_bar.fill.transparency = 0.7
        glow_bar.line.fill.background()

        # === GTM FABRIC LOGO - Top left, prominently displayed ===
        if include_logo and self.logo_path:
            logo_width = 3.5  # Larger for better visibility
            logo_height = logo_width * 0.217  # Actual logo aspect ratio

            try:
                slide.shapes.add_picture(
                    self.logo_path,
                    Inches(0.6), Inches(0.4),
                    width=Inches(logo_width)
                )
            except Exception as e:
                print(f"Warning: Could not add logo: {e}")

        # === CONTENT - Modern layout with asymmetric positioning ===
        # Title - positioned lower and left-aligned for modern look
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.3),
            Inches(8.8), Inches(1.5)
        )
        text_frame = title_box.text_frame
        text_frame.text = title
        text_frame.word_wrap = True
        text_frame.margin_left = 0
        text_frame.margin_right = 0

        para = text_frame.paragraphs[0]
        para.font.size = Pt(44)  # Larger, bolder title
        para.font.bold = True
        para.font.color.rgb = self.colors.text_primary
        para.font.name = self.font_family
        para.line_spacing = 1.1

        # Subtitle with accent color
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(3.9),
            Inches(8.8), Inches(0.6)
        )
        text_frame = subtitle_box.text_frame
        text_frame.text = subtitle
        text_frame.word_wrap = True
        text_frame.margin_left = 0
        text_frame.margin_right = 0

        para = text_frame.paragraphs[0]
        para.font.size = Pt(22)
        para.font.bold = False  # Lighter weight for contrast
        para.font.color.rgb = self.colors.accent  # Use accent color
        para.font.name = self.font_family

        # Date - minimal, bottom right
        if date:
            date_box = slide.shapes.add_textbox(
                Inches(7.5), Inches(5.0),
                Inches(2.0), Inches(0.4)
            )
            text_frame = date_box.text_frame
            text_frame.text = date
            text_frame.margin_left = 0
            text_frame.margin_right = 0

            para = text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = self.colors.text_secondary
            para.font.name = self.font_family
            para.alignment = PP_ALIGN.RIGHT

    def create_content_slide(
        self,
        prs: Presentation,
        title: str,
        content_sections: List[Dict[str, Any]],
        include_logo: bool = True,
        page_number: Optional[int] = None
    ):
        """
        Create modern content slide with card-style sections

        Args:
            prs: Presentation object
            title: Slide title
            content_sections: List of content sections with headers and text
            include_logo: Whether to include GTM Fabric logo
            page_number: Optional page number for footer

        Example content_sections:
        [
            {
                "header": "Section Title",
                "items": ["Item 1", "Item 2"],
                "style": "bullets"  # or "text"
            }
        ]
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._apply_dark_background(slide)

        # Title with accent underline
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25),
            Inches(8.5), Inches(0.5)
        )
        text_frame = title_box.text_frame
        text_frame.text = title

        para = text_frame.paragraphs[0]
        para.font.size = Pt(32)
        para.font.bold = True
        para.font.color.rgb = self.colors.text_primary
        para.font.name = self.font_family

        # Accent underline
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.72),
            Inches(2), Inches(0.03)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = self.colors.accent
        underline.line.fill.background()

        # Content card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.0),
            Inches(9), Inches(4.0)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors.card
        card.line.color.rgb = self.colors.border_card
        card.line.width = Pt(2)

        # Add content to card
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.25),
            Inches(8.4), Inches(3.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)

        first_section = True
        for section in content_sections:
            # Add spacing between sections
            if not first_section:
                p = text_frame.add_paragraph()
                p.space_before = Pt(16)
            first_section = False

            # Section header
            p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
            p.text = section.get("header", "")
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors.accent
            p.font.name = self.font_family
            p.space_after = Pt(8)

            # Section items
            style = section.get("style", "bullets")
            for item in section.get("items", []):
                p = text_frame.add_paragraph()
                if style == "bullets":
                    p.text = f"• {item}"
                    p.level = 0
                else:
                    p.text = item

                p.font.size = Pt(14)
                p.font.color.rgb = self.colors.text_primary
                p.font.name = self.font_family
                p.space_after = Pt(6)

        # Add logo aligned with title (drawn last so it appears on top)
        if include_logo:
            title_center_y = 0.25 + 0.25  # title_top + title_height/2
            self._add_logo(slide, position="top-right", size=2.0, align_center_y=title_center_y)

        # Add footer
        if page_number:
            self._add_footer(slide, "GTM Fabric", page_number)

    def create_icp_slide(
        self,
        prs: Presentation,
        icp_number: int,
        icp_title: str,
        industries: List[str],
        departments: str,
        key_roles: List[str],
        displacement_signals: List[Dict[str, Any]],
        expansion_signals: List[Dict[str, Any]],
        include_logo: bool = True,
        page_number: Optional[int] = None
    ):
        """
        Create ICP-specific slide with modern card design

        Args:
            prs: Presentation object
            icp_number: ICP number (1-4)
            icp_title: Descriptive ICP title
            industries: List of target industries
            departments: Departments description
            key_roles: List of key job roles
            displacement_signals: List of displacement signals with pain_explanation and vendor_products
            expansion_signals: List of expansion signals with readiness_explanation and vendor_products
            include_logo: Whether to include GTM Fabric logo
            page_number: Optional page number
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._apply_dark_background(slide)

        # ICP badge (modern pill design)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(0.2),
            Inches(0.9), Inches(0.3)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.colors.accent
        badge.line.fill.background()

        badge_text = badge.text_frame
        badge_text.text = f"ICP {icp_number}"
        para = badge_text.paragraphs[0]
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.color.rgb = self.colors.background
        para.font.name = self.font_family
        para.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(0.18),
            Inches(7.4), Inches(0.35)
        )
        text_frame = title_box.text_frame
        text_frame.text = icp_title

        para = text_frame.paragraphs[0]
        para.font.size = Pt(24)
        para.font.bold = True
        para.font.color.rgb = self.colors.text_primary
        para.font.name = self.font_family

        # Content card with gradient effect
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(0.65),
            Inches(9), Inches(4.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors.card
        card.line.color.rgb = self.colors.border_card
        card.line.width = Pt(1.5)

        # Content text frame
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.85),
            Inches(8.5), Inches(4.0)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        # Industries
        p = text_frame.paragraphs[0]
        p.text = "Industries"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors.accent
        p.font.name = self.font_family
        p.space_after = Pt(6)

        p = text_frame.add_paragraph()
        p.text = ", ".join(industries)
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors.text_primary
        p.font.name = self.font_family
        p.space_after = Pt(12)

        # Departments & Functions
        p = text_frame.add_paragraph()
        p.text = "Departments & Business Functions"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors.accent
        p.font.name = self.font_family
        p.space_after = Pt(6)

        p = text_frame.add_paragraph()
        p.text = departments
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors.text_primary
        p.font.name = self.font_family
        p.space_after = Pt(6)

        # Key Roles - inline format to save space
        p = text_frame.add_paragraph()
        roles_text = "Key Roles: " + ", ".join(key_roles)
        p.text = roles_text
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors.text_primary
        p.font.name = self.font_family
        p.space_after = Pt(12)

        # Technographic Fit Header
        p = text_frame.add_paragraph()
        p.text = "Technographic Fit"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors.accent
        p.font.name = self.font_family
        p.space_after = Pt(8)

        # Displacement Signals
        p = text_frame.add_paragraph()
        p.text = "Displacement/Modernization:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors.text_primary
        p.font.name = self.font_family
        p.space_after = Pt(4)

        for signal in displacement_signals:
            p = text_frame.add_paragraph()
            vendors = ", ".join(signal.get("vendor_products", []))
            p.text = f"• {signal.get('pain_explanation', '')} → {vendors}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors.text_primary
            p.font.name = self.font_family
            p.space_after = Pt(4)

        p = text_frame.add_paragraph()
        p.space_after = Pt(8)

        # Expansion Signals
        p = text_frame.add_paragraph()
        p.text = "Expansion:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors.text_primary
        p.font.name = self.font_family
        p.space_after = Pt(4)

        for signal in expansion_signals:
            p = text_frame.add_paragraph()
            vendors = ", ".join(signal.get("vendor_products", []))
            p.text = f"• {signal.get('readiness_explanation', '')} → {vendors}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors.text_primary
            p.font.name = self.font_family
            p.space_after = Pt(4)

        # Add logo aligned with title (drawn last so it appears on top)
        if include_logo:
            title_center_y = 0.18 + 0.175  # title_top + title_height/2
            self._add_logo(slide, position="top-right", size=1.8, align_center_y=title_center_y)

        # Add footer
        if page_number:
            self._add_footer(slide, "GTM Fabric", page_number)

    def create_e2e_slide(
        self,
        prs: Presentation,
        include_logo: bool = True,
        page_number: Optional[int] = None
    ):
        """Create E2E (End-to-End) Use Case Solutions slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._apply_dark_background(slide)

        # Title with accent underline
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25),
            Inches(8.5), Inches(0.5)
        )
        text_frame = title_box.text_frame
        text_frame.text = "Key Elements Required For E2E Use Case Solutions"

        para = text_frame.paragraphs[0]
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = self.colors.text_primary
        para.font.name = self.font_family

        # Accent underline
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.72),
            Inches(3), Inches(0.03)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = self.colors.accent
        underline.line.fill.background()

        # Content card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.0),
            Inches(9), Inches(4.0)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors.card
        card.line.color.rgb = self.colors.border_card
        card.line.width = Pt(2)

        # 4 circles in 2x2 grid
        circle_size = Inches(0.8)
        center_x = Inches(5)
        center_y = Inches(3)
        spacing = Inches(1.9)  # Reduced from 2.2 to bring circles closer to center

        circles = [
            {"x": center_x - spacing, "y": center_y + spacing, "text": "Data &\nInsights", "color": self.colors.accent},
            {"x": center_x + spacing, "y": center_y + spacing, "text": "Workflow &\nPlatforms", "color": self.colors.accent},
            {"x": center_x - spacing, "y": center_y - spacing, "text": "Experts &\nBlueprints", "color": self.colors.accent},
            {"x": center_x + spacing, "y": center_y - spacing, "text": "Services &\nSkills", "color": self.colors.accent},
        ]

        for circle in circles:
            # Circle shape
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                circle["x"] - circle_size/2, circle["y"] - circle_size/2,
                circle_size, circle_size
            )
            circ.fill.solid()
            circ.fill.fore_color.rgb = self.colors.card
            circ.line.color.rgb = circle["color"]
            circ.line.width = Pt(2)

            # Circle text - ensure proper centering
            text_frame = circ.text_frame
            text_frame.text = circle["text"]
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Remove any margins to ensure perfect centering
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0

            para = text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = self.colors.text_primary
            para.font.name = self.font_family
            para.alignment = PP_ALIGN.CENTER

        # Left side content (moved further left to avoid circle overlap)
        left_col_x = Inches(0.6)
        left_col_width = Inches(2.8)

        # Top left items
        top_left_items = [
            "Neutral vendor mgmt (source / vet / negotiate)",
            "Compelling value of complementary data",
            "Harmonization (integrated ontology)",
            "Flexible & scalable pricing model",
            "1st/2nd/3rd party data integration",
            "Data and insights licensing options",
        ]

        top_left_y = Inches(1.5)
        top_left_box = slide.shapes.add_textbox(
            left_col_x, top_left_y, left_col_width, Inches(1.5)
        )
        text_frame = top_left_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)

        for i, item in enumerate(top_left_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors.text_primary
            p.font.name = self.font_family
            p.space_after = Pt(2)

        # Bottom left items
        bottom_left_items = [
            "Practical examples, case studies & blueprints",
            "Proof points (economic outcomes)",
            "Success factors & watch outs",
            "Board-level credibility & fluency"
        ]

        bottom_left_y = Inches(3.5)
        bottom_left_box = slide.shapes.add_textbox(
            left_col_x, bottom_left_y, left_col_width, Inches(1.3)
        )
        text_frame = bottom_left_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)

        for i, item in enumerate(bottom_left_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors.text_primary
            p.font.name = self.font_family
            p.space_after = Pt(2)

        # Right side content (moved further right to avoid circle overlap)
        right_col_x = Inches(6.7)
        right_col_width = Inches(2.8)

        # Top right items
        top_right_items = [
            "Architecture & design",
            "Integration & customization",
            "Analytical tools & visualization",
            "AI reference framework (MCP)",
            "LLM training & governance",
            "Agentic processes & automation"
        ]

        top_right_box = slide.shapes.add_textbox(
            right_col_x, top_left_y, right_col_width, Inches(1.5)
        )
        text_frame = top_right_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)

        for i, item in enumerate(top_right_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors.text_primary
            p.font.name = self.font_family
            p.space_after = Pt(2)

        # Bottom right items
        bottom_right_items = [
            "Project management / oversight",
            "Strategic vision & design",
            "Process optimization & re-imagination",
            "Training & enablement",
            "Change management / comms"
        ]

        bottom_right_box = slide.shapes.add_textbox(
            right_col_x, bottom_left_y, right_col_width, Inches(1.3)
        )
        text_frame = bottom_right_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)

        for i, item in enumerate(bottom_right_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors.text_primary
            p.font.name = self.font_family
            p.space_after = Pt(2)

        # Draw dashed cross lines AFTER circles and text (extends to edge of content)
        from pptx.enum.dml import MSO_LINE_DASH_STYLE

        # Horizontal line through center (extends from left edge to right edge of text)
        h_line_left = left_col_x - Inches(0.1)
        h_line_right = right_col_x + right_col_width + Inches(0.1)
        line = slide.shapes.add_connector(
            1, h_line_left, center_y, h_line_right, center_y
        )
        line.line.color.rgb = RGBColor(0x5a, 0x4d, 0x6f)  # Medium purple-gray for subtle visibility
        line.line.width = Pt(1.5)  # Slightly thicker than default but not too bold
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # Vertical line through center (extends from top to bottom of text areas)
        v_line_top = top_left_y - Inches(0.1)
        v_line_bottom = bottom_left_y + Inches(1.3) + Inches(0.1)
        line = slide.shapes.add_connector(
            1, center_x, v_line_top, center_x, v_line_bottom
        )
        line.line.color.rgb = RGBColor(0x5a, 0x4d, 0x6f)  # Medium purple-gray for subtle visibility
        line.line.width = Pt(1.5)  # Slightly thicker than default but not too bold
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # Add top-right logo aligned with title (drawn after card so it's visible)
        if include_logo:
            title_center_y = 0.25 + 0.25  # title_top + title_height/2
            self._add_logo(slide, position="top-right", size=1.5, align_center_y=title_center_y)

        # Center logo (white GTM Fabric logo) - added last to be on top
        if self.logo_path:
            try:
                logo_size = Inches(2.5)
                slide.shapes.add_picture(
                    self.logo_path,
                    center_x - logo_size/2,
                    center_y - logo_size/2,
                    width=logo_size
                )
            except Exception as e:
                # Fallback to text if logo fails
                logo_box = slide.shapes.add_textbox(
                    center_x - Inches(0.75), center_y - Inches(0.2),
                    Inches(1.5), Inches(0.4)
                )
                text_frame = logo_box.text_frame
                text_frame.text = "GTM Fabric"
                para = text_frame.paragraphs[0]
                para.font.size = Pt(14)
                para.font.bold = True
                para.font.color.rgb = self.colors.accent
                para.font.name = self.font_family
                para.alignment = PP_ALIGN.CENTER
        else:
            # Fallback to text if no logo
            logo_box = slide.shapes.add_textbox(
                center_x - Inches(0.75), center_y - Inches(0.2),
                Inches(1.5), Inches(0.4)
            )
            text_frame = logo_box.text_frame
            text_frame.text = "GTM Fabric"
            para = text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = self.colors.accent
            para.font.name = self.font_family
            para.alignment = PP_ALIGN.CENTER

        # Add footer
        if page_number:
            self._add_footer(slide, "GTM Fabric", page_number)

    def create_propensity_funnel_slide(
        self,
        prs: Presentation,
        include_logo: bool = True,
        page_number: Optional[int] = None
    ):
        """Create Propensity Funnel slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._apply_dark_background(slide)

        # Title with accent underline
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25),
            Inches(8.5), Inches(0.5)
        )
        text_frame = title_box.text_frame
        text_frame.text = "Propensity Funnel"

        para = text_frame.paragraphs[0]
        para.font.size = Pt(32)
        para.font.bold = True
        para.font.color.rgb = self.colors.text_primary
        para.font.name = self.font_family

        # Accent underline
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.72),
            Inches(2), Inches(0.03)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = self.colors.accent
        underline.line.fill.background()

        # Content card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.0),
            Inches(9), Inches(4.0)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors.card
        card.line.color.rgb = self.colors.border_card
        card.line.width = Pt(2)

        # Two columns
        left_x = Inches(0.8)
        right_x = Inches(5.5)
        start_y = Inches(1.3)

        levels = [
            {"title": "Market to Account View", "desc": "TAM-SAM-SOM, Spend Potential",
             "short": "4.4M companies, $127bn annually"},
            {"title": "Account to Product View", "desc": "Competitive & Partner Install Base",
             "short": "11,759 using IBM QRadar"},
            {"title": "Product to User View", "desc": "Departments To Target, Locations",
             "short": "724 intent: Ransomware + SIEM"},
            {"title": "User to Contact Data", "desc": "Emails, Phone Numbers, Socials",
             "short": "Contacts by Dept, Location, Seniority"},
        ]

        y_pos = start_y
        for i, level in enumerate(levels):
            # Left column - Level details
            detail_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left_x, y_pos,
                Inches(4), Inches(0.7)
            )
            detail_box.fill.solid()
            detail_box.fill.fore_color.rgb = self.colors.card
            detail_box.line.color.rgb = self.colors.border_card
            detail_box.line.width = Pt(1)

            # Title
            text_frame = detail_box.text_frame
            text_frame.text = level["title"]
            text_frame.margin_left = Inches(0.15)
            text_frame.margin_top = Inches(0.1)

            para = text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = self.colors.text_primary
            para.font.name = self.font_family

            # Description
            p = text_frame.add_paragraph()
            p.text = level["desc"]
            p.font.size = Pt(9)
            p.font.color.rgb = self.colors.text_secondary
            p.font.name = self.font_family

            # Right column - Funnel visualization
            width_percent = 1.0 - (i * 0.18)
            funnel_width = Inches(3.5) * width_percent
            funnel_x = right_x + (Inches(3.5) - funnel_width) / 2

            funnel_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                funnel_x, y_pos + Inches(0.1),
                funnel_width, Inches(0.5)
            )
            funnel_box.fill.solid()
            funnel_box.fill.fore_color.rgb = self.colors.card
            funnel_box.line.color.rgb = self.colors.accent
            funnel_box.line.width = Pt(2)

            # Funnel text
            text_frame = funnel_box.text_frame
            text_frame.text = level["short"]
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.word_wrap = True  # Enable word wrapping
            # Increase margins to keep text within bounds
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)

            para = text_frame.paragraphs[0]
            # Progressively reduce font size for narrower funnels
            if i >= 3:
                para.font.size = Pt(7)  # Even smaller for level 4
            elif i >= 2:
                para.font.size = Pt(7.5)
            else:
                para.font.size = Pt(8)
            para.font.bold = True
            para.font.color.rgb = self.colors.text_primary
            para.font.name = self.font_family
            para.alignment = PP_ALIGN.CENTER
            para.line_spacing = 0.85  # Tighter line spacing for wrapped text

            y_pos += Inches(0.85)

        # Add logo aligned with title (drawn last so it appears on top)
        if include_logo:
            title_center_y = 0.25 + 0.25  # title_top + title_height/2
            self._add_logo(slide, position="top-right", size=1.5, align_center_y=title_center_y)

        # Add footer
        if page_number:
            self._add_footer(slide, "GTM Fabric", page_number)

    def save_presentation(self, prs: Presentation, output_path: str):
        """
        Save presentation to file

        Args:
            prs: Presentation object
            output_path: Output file path
        """
        prs.save(output_path)
        print(f"✓ Presentation saved to: {output_path}")


# Convenience function for quick ICP presentation generation
def generate_icp_presentation(
    campaign_info: Dict[str, str],
    icps: List[Dict[str, Any]],
    output_path: str,
    include_logo: bool = True
) -> str:
    """
    Generate complete ICP presentation from JSON data

    Args:
        campaign_info: Dict with customer_name, solution_name, campaign_date
        icps: List of ICP dicts with all required fields
        output_path: Where to save the presentation
        include_logo: Whether to include GTM Fabric logo

    Returns:
        Path to saved presentation
    """
    generator = GTMFabricPresentationGenerator()
    prs = generator.create_presentation()

    # Title slide
    generator.create_title_slide(
        prs,
        title=f"{campaign_info['customer_name']} - Ideal Customer Profiles",
        subtitle=f"{campaign_info['solution_name']} GTM Campaign",
        date=campaign_info.get('campaign_date'),
        include_logo=include_logo
    )

    # E2E slide (slide 2)
    generator.create_e2e_slide(prs, include_logo=include_logo, page_number=2)

    # Propensity Funnel slide (slide 3)
    generator.create_propensity_funnel_slide(prs, include_logo=include_logo, page_number=3)

    # ICP slides (starting from slide 4)
    for idx, icp in enumerate(icps, start=1):
        generator.create_icp_slide(
            prs,
            icp_number=icp['icp_number'],
            icp_title=icp['title'],
            industries=icp['industries'],
            departments=icp['departments_and_functions'],
            key_roles=icp['key_roles'],
            displacement_signals=icp['technographic_fit']['displacement_signals'],
            expansion_signals=icp['technographic_fit']['expansion_signals'],
            include_logo=include_logo,
            page_number=idx + 3  # Start from slide 4 (after title, E2E, and funnel)
        )

    # Save
    generator.save_presentation(prs, output_path)
    return output_path


# Example usage
if __name__ == "__main__":
    # Create example presentation
    generator = GTMFabricPresentationGenerator()
    prs = generator.create_presentation()

    # Title slide
    generator.create_title_slide(
        prs,
        title="GTM Fabric Design System",
        subtitle="Modern PowerPoint Generation",
        date="January 2025"
    )

    # Example content slide
    generator.create_content_slide(
        prs,
        title="Key Features",
        content_sections=[
            {
                "header": "Design System",
                "items": [
                    "Dark purple theme with accent colors",
                    "Modern card-based layouts",
                    "GTM Fabric branding"
                ],
                "style": "bullets"
            },
            {
                "header": "Typography",
                "items": [
                    "Arial font family",
                    "Responsive sizing",
                    "High contrast for readability"
                ],
                "style": "bullets"
            }
        ],
        page_number=2
    )

    # Save example
    output_path = "/tmp/gtm-fabric-example.pptx"
    generator.save_presentation(prs, output_path)
    print(f"\nExample presentation created: {output_path}")
